# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

GREEN = RGBColor(46, 125, 50)
DARK = RGBColor(33, 33, 33)
WHITE = RGBColor(255, 255, 255)
LIGHT_GREEN = RGBColor(200, 230, 201)
ORANGE = RGBColor(239, 108, 0)
GRAY = RGBColor(100, 100, 100)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

def tx(slide, l, t, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def mt(slide, l, t, w, h, lines, size=16, color=WHITE, bold=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf

def bar(slide, l, t, w=1.2, h=0.04):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = GREEN; s.line.fill.background()

def top_line(slide):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = GREEN; s.line.fill.background()

# S1: Title
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 1, 1.5, 11, 1.2, "元胡种植分析报告", 44, True, GREEN)
tx(s, 1, 2.8, 11, 0.8, "种植方法 · 采挖加工 · 成本测算 · 机械化方案", 22, False, RGBColor(180,180,180))
bar(s, 1, 3.6, 2)
mt(s, 1, 4, 8, 2, ["涵盖陕西汉中、河南、浙江三大产区","2024-2026年各产区调研数据","机械化耕作实践与设备厂家信息"], 16, RGBColor(160,160,160))
tx(s, 1, 6.5, 6, 0.5, "数据来源：各产区调研 · 2026年1月", 11, False, GRAY)

# S2: Contents
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, 0.5, 0.3, 12, 0.8, "目  录", 36, True, GREEN)
bar(s, 0.5, 1.0)
items = ["概述与产区格局","种植方法（整地、播种、田间管理）","采挖与初加工","成本测算与收益分析","机械化方案与设备厂家","风险提示与种植建议"]
for i, title in enumerate(items):
    y = 1.5 + i * 0.85
    tx(s, 1, y, 0.8, 0.6, f"0{i+1}", 28, True, GREEN)
    tx(s, 2, y+0.05, 8, 0.6, title, 20, False, WHITE)
    if i < len(items)-1:
        ln = s.shapes.add_shape(1, Inches(1), Inches(y+0.7), Inches(10), Inches(0.01))
        ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(60,60,60); ln.line.fill.background()

# S3: Overview
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "一、概述", 32, True, GREEN)
bar(s, 0.5, 0.9)
mt(s, 0.5, 1.3, 7, 3.5, ["元胡（延胡索）为罂粟科植物，以干燥块茎入药","功效：活血、利气、止痛","种植周期：秋播 → 次年夏收（跨年栽培）","亩产鲜品：300-400公斤（高产可达500-600公斤）","折干率：3:1（3公斤鲜品出1公斤干货）"], 17, WHITE)
tx(s, 8, 1.3, 5, 0.6, "全国三大产区格局", 20, True, GREEN)
for i, (n, sh, adv) in enumerate([("陕西汉中","核心主产区，占全国60%以上","成本最低，机械化程度高"),("河南","新兴产区，约占20%","土地劳动力有优势"),("浙江磐安","传统道地产区，不足10%","品质优、溢价25%")]):
    y = 2.1 + i * 1.6
    tx(s, 8.2, y, 4.5, 0.5, "▸ "+n, 18, True, ORANGE if i==0 else WHITE)
    tx(s, 8.5, y+0.5, 4.5, 0.4, sh, 14, False, RGBColor(180,180,180))
    tx(s, 8.5, y+0.9, 4.5, 0.4, adv, 14, False, RGBColor(160,160,160))

# S4: Planting
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "二、种植方法 — 环境要求与播前准备", 28, True, GREEN)
bar(s, 0.5, 0.9)
tx(s, 0.5, 1.3, 5, 0.5, "产地与环境要求", 20, True, LIGHT_GREEN)
mt(s, 0.5, 1.9, 5.5, 2.5, ["气候：喜温凉湿润，适宜气温5-25℃","土壤：砂质壤土，耕层≥30cm，排水好","模式：水稻+元胡水旱轮作减少病虫害"], 16, WHITE)
tx(s, 7, 1.3, 5.5, 0.5, "播前准备", 20, True, LIGHT_GREEN)
mt(s, 7, 1.9, 5.5, 3, ["基肥：农家肥1500-2000kg + 过磷酸钙40-50kg + 硫酸钾25kg/亩","选种：无病虫、饱满子元胡，千粒重约1200g","处理：25%甲霜灵800倍液浸泡10-15分钟","晾晒1-2天至表皮略有皱缩后播种"], 16, WHITE)
tx(s, 0.5, 4.5, 5, 0.5, "播种技术参数", 20, True, LIGHT_GREEN)
for i, (k, v) in enumerate([("播种时间","平坝9.20-10.10；浅山9.10-9.30"),("播种密度","行株距15x7.5cm，亩播5.8-6.0万株"),("亩用种量","75-90公斤；免耕点播约40-45公斤"),("播种深度","5-8厘米，芽向上"),("覆盖","稻草或栏肥750-1000kg覆盖")]):
    y = 5.1 + i * 0.4
    tx(s, 0.5, y, 2.5, 0.35, k, 14, True, GREEN)
    tx(s, 3, y, 4, 0.35, v, 14, False, WHITE)

# S5: Field Management
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "二、种植方法 — 田间管理", 28, True, GREEN)
bar(s, 0.5, 0.9)
tx(s, 0.5, 1.3, 6, 0.5, "施肥管理", 20, True, LIGHT_GREEN)
mt(s, 0.5, 1.9, 6, 2.5, ["基肥：农家肥1500-2000kg + 过磷酸钙40-50kg + 硫酸钾25kg/亩","腊肥（11月底-12月初）：尿素10-15kg或复合肥15kg/亩","苗肥（2月初）：尿素5-10kg/亩","叶面追肥（3-4月）：0.3%-0.4%磷酸二氢钾，7-10天一次"], 15, WHITE)
tx(s, 7, 1.3, 5.5, 0.5, "水分与除草管理", 20, True, LIGHT_GREEN)
mt(s, 7, 1.9, 5.5, 2, ["水分：怕旱怕渍，保持土壤湿润","3-4月块茎膨大期需充足水分","多雨季节及时清沟排湿","除草：芽前封闭→出苗前草甘膦→齐苗后人工"], 15, WHITE)
tx(s, 0.5, 4.5, 6, 0.5, "病虫害防治", 20, True, LIGHT_GREEN)
mt(s, 0.5, 5.1, 6, 2, ["主要病害：霜霉病（预防为主、综合防治）","水旱轮作 + 做好排水","3月中旬起，7天一次交替喷施霜脲锰锌、甲霜灵","连续2-3次，禁止使用膨大素"], 15, WHITE)

# S6: Harvest
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "三、采挖与初加工", 28, True, GREEN)
bar(s, 0.5, 0.9)
tx(s, 0.5, 1.3, 5.5, 0.5, "采挖", 20, True, LIGHT_GREEN)
mt(s, 0.5, 1.9, 5.5, 3, ["时间：4月下旬至5月中旬（茎叶枯黄时）","方法：先浅翻拣拾 → 再深翻复拣","注意：避免伤破块茎，晴天土壤干燥时进行","机械化：汉中推广，人工成本降至1500元/亩","分选：按大中小三级，有芽眼者留种"], 15, WHITE)
tx(s, 7, 1.3, 5.5, 0.5, "初加工", 20, True, LIGHT_GREEN)
mt(s, 7, 1.9, 5.5, 3.5, ["1. 清洗：箩筐搓擦表皮，洗净沥干","2. 水煮：大块4-6分钟，小块3-4分钟","  煮至切面中心呈黄色（留米粒大白点）","3. 干燥：曝晒3-4天→回潮→再晒2-3天","  阴雨天：40-60℃烘干","4. 折干率：3:1（鲜品→干货）"], 15, WHITE)

# S7: Cost
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "四、成本测算", 28, True, GREEN)
bar(s, 0.5, 0.9)
tx(s, 0.3, 1.2, 12, 0.5, "亩投入成本明细", 18, True, LIGHT_GREEN)
for i, (item, qty, price, cost) in enumerate([("种子","75-90kg","8-12元/kg","600-1080元"),("肥料","农家肥+复合肥","--","400-600元"),("农药","除草剂+杀菌剂","--","100-150元"),("人工(管理)","播种施肥等","--","800-1200元"),("采挖人工","机械/人工","1500元/亩","1500-2000元"),("初加工","清洗煮晒","约7元/kg干货","700-900元"),("土地租金","1亩","500-800元/亩","500-800元"),("其他","灌溉等","--","100-200元")]):
    y = 1.8 + i * 0.45
    tx(s, 0.5, y, 2.5, 0.4, item, 14, i==0, GREEN if i==0 else WHITE)
    tx(s, 3.2, y, 2, 0.4, qty, 13, False, RGBColor(180,180,180))
    tx(s, 5.5, y, 2, 0.4, price, 13, False, RGBColor(180,180,180))
    tx(s, 7.5, y, 2, 0.4, cost, 14, True, ORANGE)
tx(s, 0.5, 6.0, 8, 0.4, "合计：约 4,700 - 7,930 元/亩", 18, True, GREEN)
tx(s, 9, 1.2, 4, 0.5, "产出与收益（中产）", 18, True, LIGHT_GREEN)
mt(s, 9, 1.8, 4, 3, ["亩投入：5,500-6,500元","亩产值：9,300-10,640元","亩净利润：约3,000-4,000元","投入产出比：1:1.5-1.7"], 14, WHITE)

# S8: Mechanization
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "五、机械化方案", 28, True, GREEN)
bar(s, 0.5, 0.9)
tx(s, 0.5, 1.3, 6, 0.5, "智能元胡播种机", 22, True, ORANGE)
mt(s, 0.5, 1.9, 5.5, 2.5, ["研发：城固县（全国首款，6代迭代）","效率：小地块10亩/天，大地块20亩/天","节省：每亩节省人工成本约700元","厂家：城固县永生农机有限公司","电话：(0916)7366531"], 15, WHITE)
tx(s, 7, 1.3, 6, 0.5, "元胡收获机", 22, True, ORANGE)
mt(s, 7, 1.9, 5.5, 2, ["推广：陕西洋县，全县大范围应用","效率：一天可收2亩多","节省：亩均成本降低约30%","起净率：达95%以上"], 15, WHITE)
tx(s, 7, 3.8, 5.5, 0.5, "收获机厂家", 18, True, LIGHT_GREEN)
mt(s, 7, 4.3, 5.5, 2.5, ["新乡地隆药业机械 (0373-7108687)","曲阜融兴机械 (0537-4567877)","安国尚锐农机 (13722201868)"], 14, WHITE)
tx(s, 0.5, 5.3, 12, 0.5, "综合效益", 18, True, LIGHT_GREEN)
mt(s, 0.5, 5.8, 12, 0.8, ["播种省700元 + 采挖省500-600元 = 合计节约约1,200-1,300元/亩，综合降幅15-20%"], 16, ORANGE, bold=True)

# S9: Comparison
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "五、机械化方案 — 人机对比", 28, True, GREEN)
bar(s, 0.5, 0.9)
tx(s, 0.5, 1.2, 6, 0.5, "播种环节", 20, True, LIGHT_GREEN)
for i, (a, b, c) in enumerate([("对比项","人工","机械"),("效率","4-5人1天种1亩","半小时种1亩"),("成本","人工成本高","节省约700元/亩"),("用种量","浪费多","精准减少"),("工序","分步进行","一气呵成")]):
    y = 1.8 + i * 0.45
    tx(s, 0.5, y, 2, 0.4, a, 13, i==0, GREEN if i==0 else WHITE)
    tx(s, 2.8, y, 3.5, 0.4, b, 13, False, RGBColor(180,180,180) if i==0 else WHITE)
    tx(s, 6.5, y, 3.5, 0.4, c, 13, False, ORANGE if i>0 else RGBColor(180,180,180))

tx(s, 0.5, 4.3, 6, 0.5, "采挖环节", 20, True, LIGHT_GREEN)
for i, (a, b, c) in enumerate([("对比项","人工","机械"),("效率","半天挖不完半亩","一天收2亩多"),("成本","1500-2000元/亩","降低约30%"),("及时性","周期长","为晾晒争取时间")]):
    y = 4.9 + i * 0.45
    tx(s, 0.5, y, 2, 0.4, a, 13, i==0, GREEN if i==0 else WHITE)
    tx(s, 2.8, y, 3.5, 0.4, b, 13, False, RGBColor(180,180,180) if i==0 else WHITE)
    tx(s, 6.5, y, 3.5, 0.4, c, 13, False, ORANGE if i>0 else RGBColor(180,180,180))

tx(s, 10, 1.2, 3.5, 0.5, "推广现状", 18, True, LIGHT_GREEN)
mt(s, 10, 1.8, 3.5, 3.5, ["城固：500亩示范基地","洋县：全县2万余亩推广","涵盖6大环节","将制定机械化生产标准"], 14, WHITE)

# S10: Risk
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "六、风险提示与种植建议", 28, True, GREEN)
bar(s, 0.5, 0.9)
tx(s, 0.5, 1.3, 6, 0.5, "主要风险", 20, True, ORANGE)
mt(s, 0.5, 1.9, 6, 3, ["价格波动：周期性明显，2025年曾跌至36-39元/kg","病害风险：霜霉病为主要病害","气候风险：怕旱怕渍","市场风险：河南产区2025年亏损，面积缩减"], 15, WHITE)
tx(s, 7, 1.3, 5.5, 0.5, "种植建议", 20, True, LIGHT_GREEN)
mt(s, 7, 1.9, 5.5, 3.5, ["新手从5-10亩试种开始","优先汉中模式（轮作+机械化）","第一年买优质种，次年可自留","提前对接合作社或药企","推广机械化，降低人工成本40%","关注农机补贴政策"], 15, WHITE)

# S11: Contacts
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 0.5, 0.3, 10, 0.7, "设备厂家联系方式", 28, True, GREEN)
bar(s, 0.5, 0.9)
tx(s, 0.5, 1.3, 5, 0.5, "播种机", 22, True, ORANGE)
mt(s, 0.5, 1.9, 5, 2.5, ["城固县永生农机有限公司","电话：(0916)7366531","地址：陕西汉中城固县","特点：全国首款智能元胡播种机","历经6代迭代，销往8省市"], 15, WHITE)
tx(s, 6.5, 1.3, 6, 0.5, "收获机", 22, True, ORANGE)
mt(s, 6.5, 1.9, 6, 5, ["新乡地隆药业机械","电话：0373-7108687","河南新乡 · 根茎类药材收获机","","曲阜融兴机械设备","电话：0537-4567877","山东曲阜 · RX-800型，起净率98%","","安国尚锐农业机械","电话：13722201868","河北安国 · 140型滚筒式"], 14, WHITE)
tx(s, 0.5, 5.5, 12, 0.5, "建议：提前准备好种植面积、土壤类型、拖拉机马力等信息再联系厂家", 14, False, GREEN)

# S12: End
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); top_line(s)
tx(s, 1, 2, 11, 1, "谢谢观看", 48, True, GREEN, PP_ALIGN.CENTER)
bar(s, 5.5, 3.2, 2)
tx(s, 1, 3.6, 11, 0.6, "数据来源：2024-2026年各产区调研 · 2026年1月行情", 16, False, GRAY, PP_ALIGN.CENTER)
tx(s, 1, 5, 11, 0.5, "本报告仅供参考，实际种植请结合当地具体情况调整", 14, False, GRAY, PP_ALIGN.CENTER)

output_path = os.path.expanduser("~/Desktop/元胡种植分析报告.pptx")
prs.save(output_path)
print(f"DONE: {output_path}")
print(f"Total slides: {len(prs.slides)}")
